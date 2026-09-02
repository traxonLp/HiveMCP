"""Reading a template's structure so a model can fill it correctly.

This is the piece the output quality depends on. Without it a model guesses layout names
and placeholder keys; with it, it is told exactly which layouts exist, what each one can
hold, and which ``{{variables}}`` the template author left behind.

The result is written for a reader with no memory and a token budget: short field names,
no XML, and layouts described by what they *hold* rather than by their internal ids.
"""

from __future__ import annotations

import logging
import re
import zipfile
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)

PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*([A-Za-z0-9_.-]{1,64})\s*\}\}")

# An OOXML file is a zip, so an uploaded template is an untrusted archive. These bound
# the two classic decompression attacks before any parser sees the contents.
MAX_UNCOMPRESSED_BYTES = 200 * 1024 * 1024
MAX_COMPRESSION_RATIO = 200
MAX_MEMBERS = 5000

# Placeholder types that can hold slide content, as opposed to titles and furniture
# (date, footer, slide number). Used to tell a one-column layout from a two-column one.
CONTENT_PLACEHOLDER_TYPES = {"body", "object", "vertical_body", "vertical_object"}
TITLE_PLACEHOLDER_TYPES = {"title", "center_title", "vertical_title"}


def placeholder_type_name(placeholder_type: object) -> str:
    """Name a placeholder type without hardcoding numbers.

    An earlier version mapped the ECMA-376 numbering by hand and mislabelled every type,
    because python-pptx's PP_PLACEHOLDER uses its own values (TITLE=1, BODY=2,
    CENTER_TITLE=3, SUBTITLE=4, OBJECT=7...). Asking the enum for its own name cannot
    drift out of sync the way a copied table can.
    """
    name = getattr(placeholder_type, "name", None)
    if isinstance(name, str):
        return name.lower()
    return str(placeholder_type).split(" ")[0].lower()


class TemplateUnreadable(Exception):
    """The file is not a usable template."""


def assert_safe_archive(path: Path) -> None:
    """Reject zip bombs before handing the file to a parser.

    python-pptx and friends will happily expand whatever they are given, so the check has
    to happen here rather than after opening.
    """
    try:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
            if len(members) > MAX_MEMBERS:
                raise TemplateUnreadable(
                    f"archive contains {len(members)} entries, over the "
                    f"{MAX_MEMBERS} limit"
                )
            total = sum(member.file_size for member in members)
            if total > MAX_UNCOMPRESSED_BYTES:
                raise TemplateUnreadable(
                    f"archive expands to {total // 1024 // 1024} MB, over the "
                    f"{MAX_UNCOMPRESSED_BYTES // 1024 // 1024} MB limit"
                )
            for member in members:
                if member.compress_size and (
                    member.file_size / member.compress_size > MAX_COMPRESSION_RATIO
                ):
                    raise TemplateUnreadable(
                        f"entry {member.filename!r} expands "
                        f"{member.file_size // max(member.compress_size, 1)}x, which "
                        "looks like a decompression bomb"
                    )
    except zipfile.BadZipFile as exc:
        raise TemplateUnreadable("the file is not a valid Office document (not a zip)") from exc


def find_placeholders(texts: list[str]) -> list[str]:
    seen: dict[str, None] = {}
    for text in texts:
        for match in PLACEHOLDER_PATTERN.findall(text or ""):
            seen.setdefault(match, None)
    return sorted(seen)


def inspect_template(path: Path, kind: str) -> dict[str, Any]:
    # Markdown templates are plain text. Running the zip-bomb guard over one would reject
    # every single valid template, so the branch comes before it rather than inside the
    # inspector table.
    if kind == "md":
        return inspect_markdown(path)

    assert_safe_archive(path)
    inspectors = {"pptx": inspect_pptx, "docx": inspect_docx, "xlsx": inspect_xlsx}
    inspector = inspectors.get(kind)
    if inspector is None:
        raise TemplateUnreadable(f"cannot inspect a template of kind {kind!r}")
    try:
        return inspector(path)
    except TemplateUnreadable:
        raise
    except Exception as exc:  # noqa: BLE001 - a broken template must not 500
        raise TemplateUnreadable(
            f"the file could not be read as a {kind} template: {exc}"
        ) from exc


# --------------------------------------------------------------------------- #
# PowerPoint
# --------------------------------------------------------------------------- #


def inspect_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(str(path))
    texts: list[str] = []
    layouts: list[dict[str, Any]] = []

    for index, layout in enumerate(presentation.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            fmt = shape.placeholder_format
            placeholders.append(
                {
                    "idx": fmt.idx,
                    "type": placeholder_type_name(fmt.type),
                    "name": shape.name,
                }
            )
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        layouts.append(
            {
                "index": index,
                "name": layout.name,
                "placeholders": placeholders,
                # What HiveMCP's own `layout` enum would pick for this one, so the model
                # can connect the template's vocabulary to the spec's.
                "maps_to_spec_layout": _guess_spec_layout(layout.name, placeholders),
            }
        )

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)

    return {
        "kind": "pptx",
        "slide_size": {
            "width_cm": round(presentation.slide_width.cm, 1),
            "height_cm": round(presentation.slide_height.cm, 1),
        },
        "layouts": layouts,
        "example_slides": len(presentation.slides),
        "placeholders": find_placeholders(texts),
        "hint": (
            "Set slide.layout to one of the maps_to_spec_layout values. Values under "
            "'placeholders' are filled by putting them in slide.placeholders, keyed "
            "without the braces."
        ),
    }


def _guess_spec_layout(name: str, placeholders: list[dict[str, Any]]) -> str:
    """Map a template's layout to the closest value of the spec's `layout` enum.

    Name first, shape second. Layout names survive translation and renaming better than
    placeholder counts do, and a template author's "Titel und Inhalt" is a stronger
    signal than the two content placeholders it happens to contain.
    """
    lowered = (name or "").lower()
    for needle, spec_layout in (
        ("title slide", "title"),
        ("titelfolie", "title"),
        ("section", "section"),
        ("abschnitt", "section"),
        ("two content", "two_content"),
        ("zwei inhalte", "two_content"),
        ("comparison", "two_content"),
        ("vergleich", "two_content"),
        ("title and content", "title_content"),
        ("titel und inhalt", "title_content"),
        ("content with caption", "title_content"),
        ("picture with caption", "image"),
        ("bild mit untertitel", "image"),
        ("blank", "blank"),
        ("leer", "blank"),
        ("title only", "table"),
        ("nur titel", "table"),
    ):
        if needle in lowered:
            return spec_layout

    content = sum(1 for item in placeholders if item["type"] in CONTENT_PLACEHOLDER_TYPES)
    if content >= 2:
        return "two_content"
    return "title_content" if content else "blank"


# --------------------------------------------------------------------------- #
# Word
# --------------------------------------------------------------------------- #


def inspect_docx(path: Path) -> dict[str, Any]:
    document = Document(str(path))

    styles: dict[str, list[str]] = {"paragraph": [], "character": [], "table": []}
    for style in document.styles:
        bucket = {1: "paragraph", 2: "character", 3: "table"}.get(int(style.type or 0))
        if bucket:
            styles[bucket].append(style.name)

    texts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            texts.extend(cell.text for cell in row.cells)
    for section in document.sections:
        texts.extend(paragraph.text for paragraph in section.header.paragraphs)
        texts.extend(paragraph.text for paragraph in section.footer.paragraphs)

    section = document.sections[0]
    return {
        "kind": "docx",
        "page": {
            "width_cm": round(section.page_width.cm, 1),
            "height_cm": round(section.page_height.cm, 1),
            "orientation": "landscape"
            if section.page_width > section.page_height
            else "portrait",
        },
        "styles": {name: sorted(values) for name, values in styles.items()},
        # Named explicitly because the renderer falls back with a warning when one is
        # missing, and a caller reading this can see it coming.
        "list_styles_present": [
            name
            for name in ("List Bullet", "List Bullet 2", "List Bullet 3", "List Number")
            if name in styles["paragraph"]
        ],
        "has_header_or_footer": any(
            section.header.paragraphs or section.footer.paragraphs
            for section in document.sections
        ),
        "placeholders": find_placeholders(texts),
        "hint": (
            "Paragraph styles listed here can be used as-is. Values under 'placeholders' "
            "are filled via spec.placeholders, keyed without the braces."
        ),
    }


# --------------------------------------------------------------------------- #
# Excel
# --------------------------------------------------------------------------- #


def inspect_xlsx(path: Path) -> dict[str, Any]:
    workbook = load_workbook(str(path), data_only=False)
    texts: list[str] = []
    sheets: list[dict[str, Any]] = []

    for worksheet in workbook.worksheets:
        header: list[str] = []
        for row in worksheet.iter_rows(min_row=1, max_row=1, values_only=True):
            header = [str(value) for value in row if value is not None]
        # Sampled, not exhaustive: a template's variables live near the top, and a large
        # workbook would otherwise make inspection slow and the answer unreadable.
        for row in worksheet.iter_rows(min_row=1, max_row=50, values_only=True):
            texts.extend(str(value) for value in row if isinstance(value, str))
        sheets.append(
            {
                "name": worksheet.title,
                "header_row": header,
                "max_row": worksheet.max_row,
                "max_column": worksheet.max_column,
                "freeze_panes": worksheet.freeze_panes,
            }
        )

    named_ranges = []
    try:
        named_ranges = sorted(workbook.defined_names.keys())
    except Exception:  # noqa: BLE001 - openpyxl's shape here has changed across versions
        logger.debug("could not read defined names", exc_info=True)

    return {
        "kind": "xlsx",
        "sheets": sheets,
        "named_ranges": named_ranges,
        "placeholders": find_placeholders(texts),
        "hint": (
            "Sheet names are reused when a spec sheet has the same name. Values under "
            "'placeholders' were found in the first 50 rows of each sheet."
        ),
    }


# --------------------------------------------------------------------------- #
# Markdown
# --------------------------------------------------------------------------- #


def inspect_markdown(path: Path) -> dict[str, Any]:
    """Report what a Markdown template offers.

    Unlike the Office formats there is no layout to discover — a Markdown template is a
    skeleton with holes. So what matters is which holes exist, and whether the template
    says where the generated document should go.

    ``{{content}}`` is that marker. Without it the body is appended, which is a
    reasonable default but rarely what the author intended, so it is reported rather than
    left to be discovered by looking at the output.
    """
    try:
        text = path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError) as exc:
        raise TemplateUnreadable(
            "this file is not readable as UTF-8 text. A Markdown template is a plain "
            ".md file."
        ) from exc

    placeholders = [name for name in find_placeholders([text]) if name != "content"]
    has_content_marker = "{{content}}" in text

    frontmatter: list[str] = []
    if text.lstrip().startswith("---"):
        _, _, rest = text.lstrip().partition("---")
        block, sep, _ = rest.partition("\n---")
        if sep:
            frontmatter = [
                line.split(":", 1)[0].strip()
                for line in block.splitlines()
                if ":" in line and not line.startswith(" ")
            ]

    return {
        "kind": "md",
        "placeholders": placeholders,
        "has_content_marker": has_content_marker,
        "frontmatter_keys": frontmatter,
        "size_bytes": len(text.encode("utf-8")),
        "hint": (
            "Fill the placeholders through spec.placeholders, keyed without the braces. "
            + (
                "The generated document replaces {{content}}."
                if has_content_marker
                else "This template has no {{content}} marker, so the generated document "
                "is appended after it. Add {{content}} where the body belongs."
            )
        ),
    }
