"""Applying edit operations to an existing document.

Two properties are worth more than the operation set itself.

**All or nothing.** Every operation runs against an in-memory copy, and the result is
only handed back if all of them succeed. A half-applied edit is worse than a failed one:
the user gets a file that looks finished and is not, and nothing says which half landed.

**A no-op is reported, not swallowed.** ``replace_text`` for a string that does not occur
succeeds trivially and changes nothing. Left silent, the user gets their document back
seemingly edited. Every operation therefore reports what it actually touched, and a zero
count becomes a warning.
"""

from __future__ import annotations

import logging
import re
from io import BytesIO
from typing import Any

from docx import Document
from openpyxl import load_workbook
from openpyxl.utils import coordinate_to_tuple
from pptx import Presentation

from ..models import EditOp
from ..render.base import MEDIA_TYPES, RenderedFile
from ..render.xlsx import FORMULA_TRIGGERS, NUMBER_FORMATS, XlsxRenderer
from .read import DocumentUnreadable

logger = logging.getLogger(__name__)


class EditError(Exception):
    """An operation could not be applied. The message names which one."""


class Applied:
    """Collects what each operation actually changed."""

    def __init__(self) -> None:
        self.lines: list[str] = []
        self.warnings: list[str] = []

    def did(self, message: str) -> None:
        self.lines.append(message)

    def changed_nothing(self, index: int, message: str) -> None:
        self.lines.append(message)
        self.warnings.append(
            f"Operation {index} matched nothing: {message} The document is unchanged in "
            "that respect — check the exact text or position with hive_read_document."
        )


def apply_edits(
    data: bytes, kind: str, operations: list[EditOp], filename: str
) -> tuple[RenderedFile, list[str]]:
    """Apply every operation, or none. Returns the new file and a per-operation log."""
    if not operations:
        raise EditError("no operations were given, so there is nothing to change")

    editors = {
        "pptx": _edit_pptx, "docx": _edit_docx, "xlsx": _edit_xlsx, "md": _edit_md
    }
    editor = editors.get(kind)
    if editor is None:
        raise EditError(f"cannot edit a document of kind {kind!r}")

    applied = Applied()
    try:
        output = editor(data, operations, applied)
    except EditError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise EditError(f"the document could not be edited: {exc}") from exc

    return (
        RenderedFile(
            data=output,
            filename=filename,
            media_type=MEDIA_TYPES[kind],
            warnings=applied.warnings,
        ),
        applied.lines,
    )


def _wrong_kind(index: int, op: Any, kind: str) -> EditError:
    return EditError(
        f"operation {index} ({op.op}) does not apply to a {kind} file. "
        "Call hive_read_document to see what this document contains."
    )


# --------------------------------------------------------------------------- #
# PowerPoint
# --------------------------------------------------------------------------- #


def _edit_pptx(data: bytes, operations: list[EditOp], applied: Applied) -> bytes:
    presentation = Presentation(BytesIO(data))

    for index, op in enumerate(operations, start=1):
        if op.op == "replace_text":
            count = _replace_in_pptx(presentation, op.find, op.replace, op.match_case)
            message = f"replaced {op.find!r} in {count} place(s)."
            applied.did(message) if count else applied.changed_nothing(index, message)

        elif op.op == "fill_placeholders":
            total = 0
            for key, value in op.values.items():
                total += _replace_in_pptx(presentation, "{{" + key + "}}", value, True)
            message = f"filled {total} placeholder occurrence(s)."
            applied.did(message) if total else applied.changed_nothing(index, message)

        elif op.op == "delete_slide":
            _check_slide(index, op.slide, len(presentation.slides))
            _drop_slide(presentation, op.slide - 1)
            applied.did(f"deleted slide {op.slide}.")

        elif op.op == "reorder_slides":
            _reorder(presentation, op.order, index)
            applied.did(f"reordered slides to {op.order}.")

        elif op.op == "set_notes":
            _check_slide(index, op.slide, len(presentation.slides))
            slide = presentation.slides[op.slide - 1]
            slide.notes_slide.notes_text_frame.text = op.notes
            applied.did(f"set speaker notes on slide {op.slide}.")

        else:
            raise _wrong_kind(index, op, "pptx")

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _check_slide(index: int, number: int, total: int) -> None:
    if not 1 <= number <= total:
        raise EditError(
            f"operation {index}: this presentation has {total} slide(s), so there is no "
            f"slide {number}."
        )


def _replace_in_pptx(presentation: Any, find: str, replace: str, match_case: bool) -> int:
    count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                count += _replace_in_frame(shape.text_frame, find, replace, match_case)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        count += _replace_in_frame(
                            cell.text_frame, find, replace, match_case
                        )
        if slide.has_notes_slide:
            count += _replace_in_frame(
                slide.notes_slide.notes_text_frame, find, replace, match_case
            )
    return count


def _replace_in_frame(frame: Any, find: str, replace: str, match_case: bool) -> int:
    """Replace within each run, then fall back to the whole paragraph.

    Run-level first because it preserves per-run formatting. Word and PowerPoint split a
    sentence across runs for reasons of their own — a spell-check pass is enough — so a
    match that spans runs is common. Rewriting the paragraph recovers those at the cost of
    flattening its formatting to the first run's, which beats silently missing the edit.
    """
    count = 0
    for paragraph in frame.paragraphs:
        count += _replace_in_paragraph(paragraph, find, replace, match_case)
    return count


def _replace_in_paragraph(paragraph: Any, find: str, replace: str, match_case: bool) -> int:
    """Run level first, then whatever is left spanning runs.

    Run-level replacement preserves each run's own formatting. Whatever still matches
    afterwards must have been split across runs, so it is repaired by rewriting the
    paragraph into its first run — which flattens the formatting of that one paragraph,
    and beats silently missing the edit.
    """
    count = 0
    for run in paragraph.runs:
        replaced, hits = _substitute(run.text, find, replace, match_case)
        if hits:
            run.text = replaced
            count += hits

    if not paragraph.runs:
        return count

    joined = "".join(run.text for run in paragraph.runs)
    repaired, spanning = _substitute(joined, find, replace, match_case)
    if spanning:
        paragraph.runs[0].text = repaired
        for run in paragraph.runs[1:]:
            run.text = ""
        count += spanning
    return count


def _substitute(text: str, find: str, replace: str, match_case: bool) -> tuple[str, int]:
    if not text or not find:
        return text, 0
    if match_case:
        return text.replace(find, replace), text.count(find)
    pattern = re.compile(re.escape(find), re.IGNORECASE)
    return pattern.subn(replace, text)[0], len(pattern.findall(text))


def _drop_slide(presentation: Any, position: int) -> None:
    id_list = presentation.slides._sldIdLst  # noqa: SLF001
    entry = list(id_list)[position]
    presentation.part.drop_rel(entry.rId)
    id_list.remove(entry)


def _reorder(presentation: Any, order: list[int], index: int) -> None:
    total = len(presentation.slides)
    if sorted(order) != list(range(1, total + 1)):
        raise EditError(
            f"operation {index}: 'order' must list each of the {total} slide(s) exactly "
            f"once, as 1-based numbers. Got {order}."
        )
    id_list = presentation.slides._sldIdLst  # noqa: SLF001
    entries = list(id_list)
    for entry in entries:
        id_list.remove(entry)
    for number in order:
        id_list.append(entries[number - 1])


# --------------------------------------------------------------------------- #
# Word
# --------------------------------------------------------------------------- #


def _edit_docx(data: bytes, operations: list[EditOp], applied: Applied) -> bytes:
    document = Document(BytesIO(data))

    for index, op in enumerate(operations, start=1):
        if op.op == "replace_text":
            count = _replace_in_docx(document, op.find, op.replace, op.match_case)
            message = f"replaced {op.find!r} in {count} place(s)."
            applied.did(message) if count else applied.changed_nothing(index, message)

        elif op.op == "fill_placeholders":
            total = 0
            for key, value in op.values.items():
                total += _replace_in_docx(document, "{{" + key + "}}", value, True)
            message = f"filled {total} placeholder occurrence(s)."
            applied.did(message) if total else applied.changed_nothing(index, message)

        elif op.op == "set_paragraph":
            paragraphs = document.paragraphs
            if not 1 <= op.paragraph <= len(paragraphs):
                raise EditError(
                    f"operation {index}: this document has {len(paragraphs)} "
                    f"paragraph(s), so there is no paragraph {op.paragraph}."
                )
            paragraph = paragraphs[op.paragraph - 1]
            if paragraph.runs:
                paragraph.runs[0].text = op.text
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.add_run(op.text)
            applied.did(f"replaced the text of paragraph {op.paragraph}.")

        elif op.op == "append_paragraph":
            style = op.style
            if style and style not in {s.name for s in document.styles}:
                raise EditError(
                    f"operation {index}: this document defines no style {style!r}. "
                    "hive_read_document lists the ones it has."
                )
            document.add_paragraph(op.text, style=style)
            applied.did(f"appended a paragraph{f' in style {style!r}' if style else ''}.")

        else:
            raise _wrong_kind(index, op, "docx")

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _replace_in_docx(document: Any, find: str, replace: str, match_case: bool) -> int:
    count = 0

    def walk(container: Any) -> int:
        found = 0
        for paragraph in container.paragraphs:
            found += _replace_in_paragraph(paragraph, find, replace, match_case)
        for table in container.tables:
            for row in table.rows:
                for cell in row.cells:
                    found += walk(cell)
        return found

    count += walk(document)
    for section in document.sections:
        count += walk(section.header)
        count += walk(section.footer)
    return count


# --------------------------------------------------------------------------- #
# Excel
# --------------------------------------------------------------------------- #


def _edit_xlsx(data: bytes, operations: list[EditOp], applied: Applied) -> bytes:
    workbook = load_workbook(BytesIO(data))

    for index, op in enumerate(operations, start=1):
        if op.op == "set_cell":
            if op.sheet not in workbook.sheetnames:
                raise EditError(
                    f"operation {index}: this workbook has no sheet {op.sheet!r}. "
                    f"It has: {', '.join(workbook.sheetnames)}."
                )
            worksheet = workbook[op.sheet]
            try:
                row, column = coordinate_to_tuple(op.cell.upper())
            except ValueError as exc:
                raise EditError(
                    f"operation {index}: {op.cell!r} is not a cell reference. Use A1 "
                    "notation, e.g. 'B7'."
                ) from exc
            if row < 1:
                # coordinate_to_tuple happily returns row 0 for "A0"; openpyxl then
                # rejects it further down with a much less useful message.
                raise EditError(
                    f"operation {index}: rows start at 1, so {op.cell!r} does not exist."
                )
            cell = worksheet.cell(row=row, column=column)
            # Same coercion and formula neutralisation as when rendering: an edit is no
            # more trustworthy an input than a generated spec.
            cell.value = XlsxRenderer._coerce(op.value, op.type)  # noqa: SLF001
            cell.number_format = NUMBER_FORMATS[op.type]
            applied.did(f"set {op.sheet}!{op.cell.upper()}.")

        elif op.op == "replace_text":
            count = 0
            for worksheet in workbook.worksheets:
                for row in worksheet.iter_rows():
                    for cell in row:
                        if not isinstance(cell.value, str):
                            continue
                        replaced, hits = _substitute(
                            cell.value, op.find, op.replace, op.match_case
                        )
                        if hits:
                            cell.value = _neutralise(replaced)
                            count += hits
            message = f"replaced {op.find!r} in {count} cell occurrence(s)."
            applied.did(message) if count else applied.changed_nothing(index, message)

        elif op.op == "fill_placeholders":
            total = 0
            for key, value in op.values.items():
                token = "{{" + key + "}}"
                for worksheet in workbook.worksheets:
                    for row in worksheet.iter_rows():
                        for cell in row:
                            if isinstance(cell.value, str) and token in cell.value:
                                cell.value = _neutralise(cell.value.replace(token, value))
                                total += 1
            message = f"filled {total} placeholder cell(s)."
            applied.did(message) if total else applied.changed_nothing(index, message)

        else:
            raise _wrong_kind(index, op, "xlsx")

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _neutralise(text: str) -> str:
    """Stop Excel evaluating replaced text as a formula."""
    return f"'{text}" if text[:1] in FORMULA_TRIGGERS else text


__all__ = ["EditError", "apply_edits", "DocumentUnreadable"]


# --------------------------------------------------------------------------- #
# Markdown
# --------------------------------------------------------------------------- #


def _edit_md(data: bytes, operations: list[EditOp], applied: Applied) -> bytes:
    """Patch a Markdown file line by line.

    Simpler than its OOXML counterparts and deliberately so: the file is text, so there
    is no document model to preserve and no risk of a structural edit invalidating a
    relationship. What is preserved instead is everything the operations do not touch —
    including whether the file ended with a newline, because flipping that turns a
    no-op edit into a one-line diff.
    """
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise EditError("this file is not valid UTF-8 text, so it is not Markdown") from exc

    trailing_newline = text.endswith("\n")
    lines = text.splitlines()

    for index, op in enumerate(operations, start=1):
        if op.op == "set_line":
            if not 1 <= op.line <= len(lines):
                raise EditError(
                    f"operation {index}: this file has {len(lines)} line(s), so there is "
                    f"no line {op.line}."
                )
            lines[op.line - 1] = op.text
            applied.did(f"replaced line {op.line}.")

        elif op.op == "append_paragraph":
            # A blank line first, or the appended text joins the previous paragraph.
            if lines and lines[-1].strip():
                lines.append("")
            lines.append(op.text)
            applied.did("appended a paragraph.")

        elif op.op == "replace_text":
            count = 0
            for position, line in enumerate(lines):
                replaced, hits = _substitute(line, op.find, op.replace, op.match_case)
                if hits:
                    lines[position] = replaced
                    count += hits
            message = f"replaced {op.find!r} in {count} place(s)."
            applied.did(message) if count else applied.changed_nothing(index, message)

        elif op.op == "fill_placeholders":
            total = 0
            for key, value in op.values.items():
                token = "{{" + key + "}}"
                for position, line in enumerate(lines):
                    if token in line:
                        lines[position] = line.replace(token, value)
                        total += line.count(token)
            message = f"filled {total} placeholder occurrence(s)."
            applied.did(message) if total else applied.changed_nothing(index, message)

        else:
            raise _wrong_kind(index, op, "md")

    output = "\n".join(lines)
    if trailing_newline:
        output += "\n"
    return output.encode("utf-8")
