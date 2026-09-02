"""Reading an uploaded document into something a model can reason about.

The point is to give a model the coordinates it needs to write edit operations: slide
numbers, paragraph numbers, sheet names, the styles the document actually defines. Those
positions are 1-based here because that is what the edit operations take, and translating
between two numbering schemes mid-task is a reliable source of off-by-one edits.

``outline`` mode exists because a full read of a large document is mostly noise. A model
usually needs to know *where* things are, not every word.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Literal

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from ..templates.inspect import TemplateUnreadable, assert_safe_archive, find_placeholders

ReadMode = Literal["outline", "full"]

# Enough to see the shape of a sheet without turning the answer into a data dump.
SAMPLE_ROWS = 20
OUTLINE_TEXT_LIMIT = 90


class DocumentUnreadable(Exception):
    """The file could not be read as the kind it claims to be."""


def read_document(path: Path, kind: str, mode: ReadMode = "outline") -> dict[str, Any]:
    # An uploaded document is an untrusted archive, and the parsers below expand whatever
    # they are given. Same guard as for templates.
    # Markdown is text, not an archive; the zip-bomb guard would reject every valid one.
    if kind != "md":
        try:
            assert_safe_archive(path)
        except TemplateUnreadable as exc:
            raise DocumentUnreadable(str(exc)) from exc

    readers = {
        "pptx": _read_pptx, "docx": _read_docx, "xlsx": _read_xlsx, "md": _read_md
    }
    reader = readers.get(kind)
    if reader is None:
        raise DocumentUnreadable(f"cannot read a document of kind {kind!r}")
    try:
        return reader(path, mode)
    except DocumentUnreadable:
        raise
    except Exception as exc:  # noqa: BLE001 - a broken upload must not 500
        raise DocumentUnreadable(f"could not read the file as {kind}: {exc}") from exc


def _clip(text: str, mode: ReadMode) -> str:
    text = (text or "").strip()
    if mode == "full" or len(text) <= OUTLINE_TEXT_LIMIT:
        return text
    return text[:OUTLINE_TEXT_LIMIT] + "…"


# --------------------------------------------------------------------------- #


def _read_pptx(path: Path, mode: ReadMode) -> dict[str, Any]:
    presentation = Presentation(str(path))
    texts: list[str] = []
    slides: list[dict[str, Any]] = []

    for number, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title is not None else None
        body: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texts.append(shape.text_frame.text)
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    body.append(_clip(paragraph.text, mode))

        entry: dict[str, Any] = {
            "slide": number,
            "layout": slide.slide_layout.name,
            "title": _clip(title, mode) if title else None,
            "text": body,
        }
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                entry["notes"] = _clip(notes, mode)
                texts.append(notes)
        entry["has_table"] = any(shape.has_table for shape in slide.shapes)
        entry["has_chart"] = any(shape.has_chart for shape in slide.shapes)
        entry["pictures"] = sum(1 for shape in slide.shapes if shape.shape_type == 13)
        slides.append(entry)

    return {
        "kind": "pptx",
        "slide_count": len(slides),
        "slides": slides,
        "placeholders": find_placeholders(texts),
        "hint": (
            "Slide numbers are 1-based and are what delete_slide, reorder_slides and "
            "set_notes expect."
        ),
    }


def _read_docx(path: Path, mode: ReadMode) -> dict[str, Any]:
    document = Document(str(path))
    texts: list[str] = []
    paragraphs: list[dict[str, Any]] = []

    for number, paragraph in enumerate(document.paragraphs, start=1):
        texts.append(paragraph.text)
        if mode == "outline" and not paragraph.text.strip():
            # Empty paragraphs are numbered but not worth listing; the number is kept so
            # set_paragraph still lines up with the document.
            continue
        paragraphs.append(
            {
                "paragraph": number,
                "style": paragraph.style.name if paragraph.style else None,
                "text": _clip(paragraph.text, mode),
            }
        )

    tables = []
    for index, table in enumerate(document.tables, start=1):
        for row in table.rows:
            texts.extend(cell.text for cell in row.cells)
        tables.append(
            {
                "table": index,
                "rows": len(table.rows),
                "columns": len(table.columns),
                "header": [cell.text for cell in table.rows[0].cells] if table.rows else [],
            }
        )

    for section in document.sections:
        texts.extend(p.text for p in section.header.paragraphs)
        texts.extend(p.text for p in section.footer.paragraphs)

    return {
        "kind": "docx",
        "paragraph_count": len(document.paragraphs),
        "paragraphs": paragraphs,
        "tables": tables,
        "styles_available": sorted(
            {style.name for style in document.styles if style.type == 1}
        ),
        "placeholders": find_placeholders(texts),
        "hint": (
            "Paragraph numbers are 1-based and count every paragraph including empty "
            "ones, which is what set_paragraph expects. In outline mode empty paragraphs "
            "are omitted from this list but still counted."
        ),
    }


def _read_xlsx(path: Path, mode: ReadMode) -> dict[str, Any]:
    workbook = load_workbook(str(path), data_only=False)
    texts: list[str] = []
    sheets: list[dict[str, Any]] = []

    limit = None if mode == "full" else SAMPLE_ROWS
    for worksheet in workbook.worksheets:
        rows: list[list[Any]] = []
        for row in worksheet.iter_rows(max_row=limit, values_only=True):
            rows.append([value for value in row])
            texts.extend(str(value) for value in row if isinstance(value, str))

        sheets.append(
            {
                "name": worksheet.title,
                "max_row": worksheet.max_row,
                "max_column": worksheet.max_column,
                "header_row": [str(v) for v in rows[0]] if rows and rows[0] else [],
                "rows": rows[1:] if mode == "full" else rows[1:SAMPLE_ROWS],
                "truncated": mode != "full" and worksheet.max_row > SAMPLE_ROWS,
            }
        )

    return {
        "kind": "xlsx",
        "sheets": sheets,
        "placeholders": find_placeholders(texts),
        "hint": "set_cell takes a sheet name and A1 notation, e.g. sheet='Data', cell='B7'.",
    }


def _read_md(path: Path, mode: ReadMode) -> dict[str, Any]:
    """Report a Markdown file as numbered lines and a heading outline.

    Lines, not paragraphs. Markdown's unit of address is the line — that is what an editor
    shows and what a diff speaks — and set_line is far easier for a model to target
    correctly than a paragraph index it would have to derive from blank-line grouping.

    Fenced code is tracked so a `#` inside a shell snippet is not reported as a heading.
    """
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()

    headings: list[dict[str, Any]] = []
    in_fence = False
    fence: str | None = None
    for number, line in enumerate(lines, start=1):
        stripped = line.strip()
        if in_fence:
            if fence and stripped.startswith(fence):
                in_fence, fence = False, None
            continue
        if stripped.startswith("```") or stripped.startswith("~~~"):
            in_fence, fence = True, stripped[:3]
            continue
        if stripped.startswith("#"):
            marker, _, title = stripped.partition(" ")
            if set(marker) == {"#"} and title.strip():
                headings.append(
                    {"line": number, "level": len(marker), "text": _clip(title.strip(), mode)}
                )

    listed = [
        {"line": number, "text": _clip(line, mode)}
        for number, line in enumerate(lines, start=1)
        if mode == "full" or line.strip()
    ]

    return {
        "kind": "md",
        "line_count": len(lines),
        "lines": listed,
        "headings": headings,
        "placeholders": find_placeholders([text]),
        "hint": (
            "Line numbers are 1-based and count every line including blank ones, which "
            "is what set_line and delete_lines expect. In outline mode blank lines are "
            "omitted from this list but still counted."
        ),
    }
