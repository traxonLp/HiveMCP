"""Render a DeckSpec into a .pptx file."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt

from .base import MEDIA_TYPES, RenderedFile, RenderError
from .theme import (
    ImageResolver,
    RenderWarnings,
    check_font,
    decode_image,
    image_size,
    normalize_hex,
    safe_filename,
)

# Aliases are matched against the layout *names* in the template, case-insensitively.
# Matching by name first is what makes custom corporate templates work: their layout
# order is arbitrary, but the names are usually recognisable. The index is only a
# fallback for the stock python-pptx template.
LAYOUT_ALIASES: dict[str, tuple[tuple[str, ...], int]] = {
    "title": (("title slide", "titelfolie", "title"), 0),
    "title_content": (("title and content", "titel und inhalt", "content"), 1),
    "section": (("section header", "abschnitt", "section"), 2),
    "two_content": (("two content", "zwei inhalte", "comparison", "vergleich"), 3),
    # Images, tables and charts are placed as standalone shapes at known coordinates
    # rather than into a layout placeholder. Placeholder geometry varies wildly between
    # templates, and a predictable position beats an inherited one that may be tiny.
    "image": (("title only", "nur titel"), 5),
    "table": (("title only", "nur titel"), 5),
    "chart": (("title only", "nur titel"), 5),
    "blank": (("blank", "leer"), 6),
}

CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "area": XL_CHART_TYPE.AREA,
}

# Body placeholder types, per ECMA-376. 1 = BODY, 7 = OBJECT, 2 = SUBTITLE.
BODY_PLACEHOLDER_TYPES = {1, 7}


class PptxRenderer:
    def __init__(
        self,
        options: Any,
        image_resolver: ImageResolver | None = None,
        template_path: Path | None = None,
    ) -> None:
        self.options = options
        self.image_resolver = image_resolver
        self.template_path = template_path
        self.warnings = RenderWarnings()
        self.font = check_font(
            getattr(options, "font_family", None),
            self.warnings,
            getattr(options, "language", "en"),
        )

    # ---------------------------------------------------------------- public

    def render(self, spec: Any) -> RenderedFile:
        prs = self._open_presentation()

        for index, slide in enumerate(spec.slides):
            try:
                self._add_slide(prs, slide)
            except Exception as exc:  # noqa: BLE001 - one bad slide must name itself
                # Deliberately re-wrapped even when it is already a RenderError: the
                # inner message says what is wrong, only this frame knows *where*, and
                # the model needs both to fix the spec on the next turn.
                raise RenderError(
                    f"slide {index + 1} ({slide.layout}) could not be rendered: {exc}"
                ) from exc

        self._apply_typography(prs)

        buffer = BytesIO()
        prs.save(buffer)
        return RenderedFile(
            data=buffer.getvalue(),
            filename=safe_filename(
                getattr(self.options, "filename", None), spec.title, "pptx"
            ),
            media_type=MEDIA_TYPES["pptx"],
            warnings=list(self.warnings),
            slide_count=len(prs.slides),
        )

    # --------------------------------------------------------------- private

    def _open_presentation(self) -> Presentation:
        if self.template_path is None:
            return Presentation()
        try:
            prs = Presentation(str(self.template_path))
        except Exception as exc:  # noqa: BLE001
            raise RenderError(
                f"template {self.template_path.name!r} could not be opened as a "
                f"PowerPoint file: {exc}"
            ) from exc
        # A .potx used as a starting point may carry example slides. Callers expect a
        # template to contribute layouts and theme, not content.
        self._remove_all_slides(prs)
        return prs

    @staticmethod
    def _remove_all_slides(prs: Presentation) -> None:
        id_list = prs.slides._sldIdLst  # noqa: SLF001
        for slide_id in list(id_list):
            prs.part.drop_rel(slide_id.rId)
            id_list.remove(slide_id)

    def _resolve_layout(self, prs: Presentation, layout_name: str):
        aliases, fallback_index = LAYOUT_ALIASES.get(layout_name, ((), 6))
        by_name = {layout.name.strip().lower(): layout for layout in prs.slide_layouts}
        for alias in aliases:
            if alias in by_name:
                return by_name[alias]
        for alias in aliases:
            for name, layout in by_name.items():
                if alias in name:
                    return layout
        if fallback_index < len(prs.slide_layouts):
            return prs.slide_layouts[fallback_index]
        self.warnings.add(
            f"Template has no layout matching {layout_name!r}; used the first available one."
        )
        return prs.slide_layouts[0]

    def _add_slide(self, prs: Presentation, spec: Any) -> None:
        layout = self._resolve_layout(prs, spec.layout)
        slide = prs.slides.add_slide(layout)

        if spec.title is not None and slide.shapes.title is not None:
            slide.shapes.title.text = spec.title

        handler = getattr(self, f"_fill_{spec.layout}", self._fill_title_content)
        handler(prs, slide, spec)

        if spec.placeholders:
            self._fill_named_placeholders(slide, spec.placeholders)
        if spec.notes:
            # Written whenever the caller supplied them. RenderOptions.include_notes is
            # a hint for brief-mode generation, not a gate on explicit content.
            slide.notes_slide.notes_text_frame.text = spec.notes

        self._drop_empty_placeholders(slide)

    # -- layout handlers ---------------------------------------------------

    def _fill_title(self, prs: Presentation, slide: Any, spec: Any) -> None:
        subtitle = spec.subtitle or spec.body
        if not subtitle:
            return
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx != 0:
                placeholder.text_frame.text = subtitle
                return

    _fill_section = _fill_title

    def _fill_title_content(self, prs: Presentation, slide: Any, spec: Any) -> None:
        body = self._first_body_placeholder(slide)
        if body is None:
            body = self._add_textbox(prs, slide)
        if spec.bullets:
            self._write_bullets(body.text_frame, spec.bullets)
        elif spec.body:
            body.text_frame.text = spec.body
        if spec.table is not None:
            self._add_table(prs, slide, spec.table, top_cm=10.0)
        if spec.chart is not None:
            self._add_chart(prs, slide, spec.chart, top_cm=10.0)
        if spec.image is not None:
            self._add_image(prs, slide, spec.image)

    def _fill_two_content(self, prs: Presentation, slide: Any, spec: Any) -> None:
        bodies = self._body_placeholders(slide)
        if len(bodies) < 2:
            self.warnings.add(
                "Layout 'two_content' has fewer than two content placeholders in this "
                "template; both columns were merged into one."
            )
            self._fill_title_content(prs, slide, spec)
            return
        self._write_bullets(bodies[0].text_frame, spec.bullets)
        self._write_bullets(bodies[1].text_frame, spec.bullets_right)

    def _fill_image(self, prs: Presentation, slide: Any, spec: Any) -> None:
        if spec.image is None:
            raise RenderError("layout 'image' requires an image")
        self._add_image(prs, slide, spec.image)
        if spec.bullets or spec.body:
            body = self._first_body_placeholder(slide)
            if body is not None:
                if spec.bullets:
                    self._write_bullets(body.text_frame, spec.bullets)
                else:
                    body.text_frame.text = spec.body

    def _fill_table(self, prs: Presentation, slide: Any, spec: Any) -> None:
        if spec.table is None:
            raise RenderError("layout 'table' requires a table")
        self._add_table(prs, slide, spec.table)

    def _fill_chart(self, prs: Presentation, slide: Any, spec: Any) -> None:
        if spec.chart is None:
            raise RenderError("layout 'chart' requires a chart")
        self._add_chart(prs, slide, spec.chart)

    def _fill_blank(self, prs: Presentation, slide: Any, spec: Any) -> None:
        if spec.bullets or spec.body:
            box = self._add_textbox(prs, slide)
            if spec.bullets:
                self._write_bullets(box.text_frame, spec.bullets)
            else:
                box.text_frame.text = spec.body
        if spec.image is not None:
            self._add_image(prs, slide, spec.image)

    # -- building blocks ---------------------------------------------------

    @staticmethod
    def _body_placeholders(slide: Any) -> list[Any]:
        found = [
            placeholder
            for placeholder in slide.placeholders
            if placeholder.placeholder_format.idx != 0
            and placeholder.placeholder_format.type in BODY_PLACEHOLDER_TYPES
        ]
        return sorted(found, key=lambda shape: (shape.top or 0, shape.left or 0))

    def _first_body_placeholder(self, slide: Any) -> Any | None:
        bodies = self._body_placeholders(slide)
        return bodies[0] if bodies else None

    @staticmethod
    def _add_textbox(prs: Presentation, slide: Any) -> Any:
        return slide.shapes.add_textbox(
            Cm(2.5), Cm(4.0), prs.slide_width - Cm(5.0), prs.slide_height - Cm(6.0)
        )

    def _write_bullets(self, text_frame: Any, bullets: list[Any], level: int = 0) -> None:
        if not bullets:
            return
        text_frame.word_wrap = True
        for bullet in bullets:
            if level == 0 and not text_frame.paragraphs[0].runs and not text_frame.text:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.level = min(level, 4)
            run = paragraph.add_run()
            run.text = bullet.text
            run.font.bold = bullet.bold or None
            run.font.italic = bullet.italic or None
            if bullet.children:
                self._write_bullets(text_frame, bullet.children, level + 1)

    def _add_table(
        self, prs: Presentation, slide: Any, data: Any, top_cm: float = 4.0
    ) -> None:
        rows = len(data.rows) + 1
        cols = len(data.headers)
        left, top = Cm(2.0), Cm(top_cm)
        width = prs.slide_width - Cm(4.0)
        height = Cm(min(1.0 * rows, 12.0))
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for index, header in enumerate(data.headers):
            cell = table.cell(0, index)
            cell.text = header
            if data.header_bold:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True

        for row_index, row in enumerate(data.rows, start=1):
            for col_index, value in enumerate(row):
                table.cell(row_index, col_index).text = str(value)

        if data.column_widths_cm:
            for index, width_cm in enumerate(data.column_widths_cm):
                table.columns[index].width = Cm(width_cm)

    def _add_chart(
        self, prs: Presentation, slide: Any, data: Any, top_cm: float = 4.0
    ) -> None:
        chart_data = CategoryChartData()
        chart_data.categories = data.categories
        for series in data.series:
            chart_data.add_series(series.name, series.values)

        graphic_frame = slide.shapes.add_chart(
            CHART_TYPES[data.chart_type],
            Cm(2.0),
            Cm(top_cm),
            prs.slide_width - Cm(4.0),
            Cm(min(11.0, prs.slide_height.cm - top_cm - 1.5)),
            chart_data,
        )
        chart = graphic_frame.chart
        if data.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = data.title
        chart.has_legend = len(data.series) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False

    def _add_image(self, prs: Presentation, slide: Any, image: Any) -> None:
        stream = decode_image(image, self.image_resolver)
        kwargs = image_size(image)
        if not kwargs:
            kwargs["width"] = prs.slide_width - Cm(6.0)
        picture = slide.shapes.add_picture(stream, Cm(3.0), Cm(4.5), **kwargs)
        if getattr(image, "alt_text", None):
            picture._element._nvXxPr.cNvPr.set("descr", image.alt_text)  # noqa: SLF001

    def _fill_named_placeholders(self, slide: Any, values: dict[str, str]) -> None:
        """Replace ``{{name}}`` tokens that the template author put in the layout."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in values.items():
                        token = "{{" + key + "}}"
                        if token in run.text:
                            run.text = run.text.replace(token, value)

    @staticmethod
    def _drop_empty_placeholders(slide: Any) -> None:
        """Remove placeholders nothing was written into.

        Left in place, empty text placeholders render as 'Click to add text' prompts and
        empty picture placeholders as grey icon boxes. Neither prints, but both make an
        otherwise finished deck look unfinished when the user opens it to review.
        Content is always added as standalone shapes, so any surviving placeholder is
        by definition unused.
        """
        for shape in list(slide.placeholders):
            is_empty_text = shape.has_text_frame and not shape.text_frame.text.strip()
            if is_empty_text or not shape.has_text_frame:
                shape._element.getparent().remove(shape._element)  # noqa: SLF001

    @staticmethod
    def _is_title(shape: Any) -> bool:
        return bool(shape.is_placeholder and shape.placeholder_format.idx == 0)

    def _apply_typography(self, prs: Presentation) -> None:
        base_size = getattr(self.options, "font_size_base", None)
        accent = normalize_hex((getattr(self.options, "theme_colors", None) or {}).get("text"))
        if not self.font and not base_size and not accent:
            return

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                is_title = self._is_title(shape)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if self.font:
                            run.font.name = self.font
                        if base_size and not is_title:
                            # Nested bullets step down 2 pt per level, floored at 8 pt.
                            run.font.size = Pt(max(8, base_size - 2 * paragraph.level))
                        if accent and not is_title:
                            run.font.color.rgb = RGBColor.from_string(accent)


def render_presentation(
    spec: Any,
    options: Any,
    image_resolver: ImageResolver | None = None,
    template_path: Path | None = None,
) -> RenderedFile:
    return PptxRenderer(options, image_resolver, template_path).render(spec)
