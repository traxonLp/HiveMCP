"""The admin-curated template pool.

Two things carry the weight here: only administrators may write, and a caller-supplied
template id must never be able to walk out of the templates root.
"""

from __future__ import annotations

import json
import os
import zipfile
from dataclasses import dataclass
from pathlib import Path

import pytest

from hivemcp.core.models import DeckSpec, DocSpec, Heading, RenderOptions, Slide
from hivemcp.core.render.docx import render_document
from hivemcp.core.render.pptx import render_presentation
from hivemcp.core.templates.inspect import (
    TemplateUnreadable,
    assert_safe_archive,
    find_placeholders,
    inspect_template,
    placeholder_type_name,
)
from hivemcp.core.templates.store import (
    NotPermitted,
    TemplateError,
    TemplateStore,
    kind_for,
    slugify,
)


@dataclass(frozen=True)
class FakeIdentity:
    user_id: str
    is_admin: bool


@pytest.fixture
def admin() -> FakeIdentity:
    return FakeIdentity(user_id="admin-1", is_admin=True)


@pytest.fixture
def member() -> FakeIdentity:
    return FakeIdentity(user_id="user-1", is_admin=False)


@pytest.fixture
def store(tmp_path: Path) -> TemplateStore:
    root = tmp_path / "templates"
    root.mkdir()
    return TemplateStore(root)


@pytest.fixture
def deck_file(options: RenderOptions) -> bytes:
    spec = DeckSpec(
        title="Corporate",
        slides=[Slide(layout="title", title="{{titel}}", subtitle="{{untertitel}}")],
    )
    return render_presentation(spec, options).data


@pytest.fixture
def doc_file(options: RenderOptions) -> bytes:
    spec = DocSpec(title="Brief", blocks=[Heading(text="{{empfaenger}}", level=1)])
    return render_document(spec, options).data


# --------------------------------------------------------------------------- #
# Permissions
# --------------------------------------------------------------------------- #


def test_admin_can_add_a_template(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    meta = store.put(deck_file, name="Corporate Deck", filename="deck.pptx", identity=admin)
    assert meta.template_id == "corporate-deck"
    assert meta.kind == "pptx"
    assert meta.created_by == "admin-1"
    assert meta.size_bytes == len(deck_file)


def test_member_cannot_add_a_template(
    store: TemplateStore, member: FakeIdentity, deck_file: bytes
) -> None:
    with pytest.raises(NotPermitted) as caught:
        store.put(deck_file, name="Meins", filename="deck.pptx", identity=member)
    # The message has to tell them what they *can* do, or the model just retries.
    assert "hive_list_templates" in str(caught.value)


def test_member_cannot_delete_a_template(
    store: TemplateStore, admin: FakeIdentity, member: FakeIdentity, deck_file: bytes
) -> None:
    meta = store.put(deck_file, name="Corporate", filename="d.pptx", identity=admin)
    with pytest.raises(NotPermitted):
        store.delete(meta.template_id, member)
    assert store.get(meta.template_id)


def test_everyone_can_read_the_shared_pool(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    """Reads are deliberately unauthenticated at this layer: the pool is shared."""
    store.put(deck_file, name="Corporate", filename="d.pptx", identity=admin)
    assert len(store.list()) == 1
    assert store.get("corporate").meta.name == "Corporate"


# --------------------------------------------------------------------------- #
# Ids and paths
# --------------------------------------------------------------------------- #


@pytest.mark.parametrize(
    ("name", "expected"),
    [
        ("Corporate Deck 2026", "corporate-deck-2026"),
        ("Angebot für Kunden", "angebot-fur-kunden"),
        ("  spaces  ", "spaces"),
        ("!!!", "template"),
        ("Ümläute & Co.", "umlaute-co"),
    ],
)
def test_slugify(name: str, expected: str) -> None:
    assert slugify(name) == expected


def test_colliding_names_get_distinct_ids(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    first = store.put(deck_file, name="Deck", filename="a.pptx", identity=admin)
    second = store.put(deck_file, name="Deck", filename="b.pptx", identity=admin)
    assert first.template_id != second.template_id
    assert {first.template_id, second.template_id} == {"deck", "deck-2"}


@pytest.mark.parametrize(
    "template_id",
    ["../etc/passwd", "..", "a/b", "/absolute", "Uppercase", "", "-leading", "x" * 100],
)
def test_malformed_ids_are_rejected_on_shape(
    store: TemplateStore, template_id: str
) -> None:
    """Rejected, not sanitised — sanitising invites the next bypass."""
    with pytest.raises(TemplateError):
        store.get(template_id)


def test_unknown_id_says_how_to_find_the_real_ones(store: TemplateStore) -> None:
    with pytest.raises(TemplateError) as caught:
        store.get("gibt-es-nicht")
    assert "hive_list_templates" in str(caught.value)


@pytest.mark.parametrize(
    ("filename", "kind"),
    [
        ("a.pptx", "pptx"),
        ("a.potx", "pptx"),
        ("a.docx", "docx"),
        ("a.dotx", "docx"),
        ("a.xlsx", "xlsx"),
        ("a.xltx", "xlsx"),
        ("a.PPTX", "pptx"),
    ],
)
def test_kind_for_accepts_documents_and_template_variants(
    filename: str, kind: str
) -> None:
    assert kind_for(filename) == kind


@pytest.mark.parametrize("filename", ["a.pdf", "a.key", "noextension", "a.txt"])
def test_kind_for_rejects_everything_else(filename: str) -> None:
    with pytest.raises(TemplateError):
        kind_for(filename)


# --------------------------------------------------------------------------- #
# Listing and metadata
# --------------------------------------------------------------------------- #


def test_list_filters_by_kind(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes, doc_file: bytes
) -> None:
    store.put(deck_file, name="Deck", filename="d.pptx", identity=admin)
    store.put(doc_file, name="Brief", filename="b.docx", identity=admin)
    assert [m.name for m in store.list("pptx")] == ["Deck"]
    assert [m.name for m in store.list("docx")] == ["Brief"]
    assert len(store.list()) == 2


def test_list_is_sorted_by_name(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    for name in ("Zebra", "alpha", "Mitte"):
        store.put(deck_file, name=name, filename="d.pptx", identity=admin)
    assert [m.name for m in store.list()] == ["alpha", "Mitte", "Zebra"]


def test_list_on_an_empty_or_missing_root(tmp_path: Path) -> None:
    assert TemplateStore(tmp_path / "does-not-exist").list() == []


def test_metadata_survives_a_new_store_instance(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    """No index database, so the meta.json beside the file is the only record."""
    meta = store.put(
        deck_file,
        name="Corporate",
        filename="d.pptx",
        identity=admin,
        description="Der Standard",
    )
    reloaded = TemplateStore(store.root).get(meta.template_id).meta
    assert reloaded.description == "Der Standard"
    assert reloaded.created_at == meta.created_at


def test_legacy_owner_id_metadata_is_still_readable(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    """An existing volume predates the rename to created_by and must keep working."""
    meta = store.put(deck_file, name="Alt", filename="d.pptx", identity=admin)
    path = store.root / meta.template_id / "meta.json"
    raw = json.loads(path.read_text())
    raw["owner_id"] = raw.pop("created_by")
    path.write_text(json.dumps(raw))

    assert TemplateStore(store.root).get(meta.template_id).meta.created_by == "admin-1"


def test_unreadable_metadata_is_skipped_rather_than_fatal(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    store.put(deck_file, name="Gut", filename="d.pptx", identity=admin)
    broken = store.root / "kaputt"
    broken.mkdir()
    (broken / "meta.json").write_text("{ not json")
    assert [m.name for m in store.list()] == ["Gut"]


def test_cache_is_invalidated_when_the_file_changes(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    """Several replicas share one volume, so mtime is the only coordination there is."""
    meta = store.put(deck_file, name="Erst", filename="d.pptx", identity=admin)
    assert store.get(meta.template_id).meta.name == "Erst"

    path = store.root / meta.template_id / "meta.json"
    raw = json.loads(path.read_text())
    raw["name"] = "Danach"
    path.write_text(json.dumps(raw))
    stat = path.stat()
    os.utime(path, (stat.st_atime, stat.st_mtime + 10))

    assert store.get(meta.template_id).meta.name == "Danach"


def test_delete_removes_the_directory(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    meta = store.put(deck_file, name="Weg", filename="d.pptx", identity=admin)
    store.delete(meta.template_id, admin)
    assert not (store.root / meta.template_id).exists()
    assert store.list() == []


def test_metadata_without_its_file_reports_corruption(
    store: TemplateStore, admin: FakeIdentity, deck_file: bytes
) -> None:
    meta = store.put(deck_file, name="Halb", filename="d.pptx", identity=admin)
    (store.root / meta.template_id / meta.filename).unlink()
    with pytest.raises(TemplateError) as caught:
        store.get(meta.template_id)
    assert "corrupt" in str(caught.value)


@pytest.mark.skipif(
    hasattr(os, "geteuid") and os.geteuid() == 0,
    reason="root ignores the permission bits this test relies on",
)
def test_an_unwritable_root_explains_itself(
    tmp_path: Path, admin: FakeIdentity, deck_file: bytes
) -> None:
    """The volume being unwritable is a deployment problem, not a 500.

    This is the failure the /templates mount-point ownership bug produced: an empty
    named volume mounted at a path the image never created belongs to root, the
    container runs as uid 10001, and the first upload dies on write.
    """
    root = tmp_path / "ro"
    root.mkdir()
    root.chmod(0o500)
    try:
        with pytest.raises(TemplateError) as caught:
            TemplateStore(root).put(
                deck_file, name="X", filename="d.pptx", identity=admin
            )
        assert "writable" in str(caught.value)
    finally:
        root.chmod(0o700)


# --------------------------------------------------------------------------- #
# Inspection
# --------------------------------------------------------------------------- #


def test_inspect_pptx_reports_layouts_and_placeholders(
    tmp_path: Path, deck_file: bytes
) -> None:
    path = tmp_path / "t.pptx"
    path.write_bytes(deck_file)
    report = inspect_template(path, "pptx")
    assert report["kind"] == "pptx"
    assert report["layouts"]
    assert all("name" in layout for layout in report["layouts"])
    assert "titel" in report["placeholders"]


def test_inspect_maps_layouts_to_spec_layout_names(
    tmp_path: Path, deck_file: bytes
) -> None:
    """The mapping is the reason to call this at all: it tells the model what to pass."""
    path = tmp_path / "t.pptx"
    path.write_bytes(deck_file)
    valid = {"title", "title_content", "two_content", "section", "image", "table", "chart", "blank"}
    guessed = {layout.get("spec_layout") for layout in inspect_template(path, "pptx")["layouts"]}
    assert guessed - {None} <= valid


def test_inspect_docx_reports_styles(tmp_path: Path, doc_file: bytes) -> None:
    path = tmp_path / "t.docx"
    path.write_bytes(doc_file)
    report = inspect_template(path, "docx")
    assert report["kind"] == "docx"
    assert report["styles"]
    assert "empfaenger" in report["placeholders"]


def test_placeholder_type_name_uses_the_enums_own_name() -> None:
    """Hardcoding this table once mapped "Title and Content" onto two_content."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    assert placeholder_type_name(PP_PLACEHOLDER.TITLE).lower().startswith("title")
    assert placeholder_type_name(PP_PLACEHOLDER.BODY).lower().startswith("body")
    assert placeholder_type_name(None)


@pytest.mark.parametrize(
    ("texts", "expected"),
    [
        (["{{a}} und {{b}}"], ["a", "b"]),
        (["{{ mit_leerzeichen }}"], ["mit_leerzeichen"]),
        (["{{a}}", "{{a}}"], ["a"]),
        (["nichts hier"], []),
        ([""], []),
    ],
)
def test_find_placeholders(texts: list[str], expected: list[str]) -> None:
    assert find_placeholders(texts) == expected


def test_a_non_zip_is_rejected(tmp_path: Path) -> None:
    path = tmp_path / "x.pptx"
    path.write_bytes(b"definitely not a zip")
    with pytest.raises(TemplateUnreadable):
        assert_safe_archive(path)


def test_a_decompression_bomb_is_rejected(tmp_path: Path) -> None:
    """Uploads are untrusted, and the parsers expand whatever they are handed."""
    path = tmp_path / "bomb.pptx"
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("payload", b"\0" * (60 * 1024 * 1024))
    with pytest.raises(TemplateUnreadable):
        assert_safe_archive(path)


def test_a_real_document_passes_the_archive_check(
    tmp_path: Path, deck_file: bytes
) -> None:
    path = tmp_path / "ok.pptx"
    path.write_bytes(deck_file)
    assert_safe_archive(path)


def test_inspecting_an_unknown_kind_is_refused(tmp_path: Path, deck_file: bytes) -> None:
    path = tmp_path / "t.pptx"
    path.write_bytes(deck_file)
    with pytest.raises(TemplateUnreadable):
        inspect_template(path, "pdf")
