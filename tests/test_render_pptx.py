"""PowerPoint rendering.

These assert against the *output file*, not against the renderer's internals: every test
opens the produced .pptx with python-pptx and checks what a viewer would see. A test that
only checked ``render_presentation`` returned bytes would have passed throughout the two
layout bugs this file now guards.
"""

from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation
from pptx.util import Emu

from hivemcp.core.models import (
    Bullet,
    ChartData,
    ChartSeries,
    DeckSpec,
    ImageRef,
    RenderOptions,
    Slide,
    TableData,
)
from hivemcp.core.render.base import MEDIA_TYPES, RenderError
from hivemcp.core.render.pptx import render_presentation


def opened(rendered) -> Presentation:
    return Presentation(BytesIO(rendered.data))


def all_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


# --------------------------------------------------------------------------- #
# Structure
# --------------------------------------------------------------------------- #


def test_renders_every_slide_in_order(deck: DeckSpec, options: RenderOptions) -> None:
    presentation = opened(render_presentation(deck, options))
    assert len(presentation.slides) == len(deck.slides)
    assert "Q3 2026" in all_text(presentation.slides[0])
    assert "Ergebnisse" in all_text(presentation.slides[1])


def test_reports_slide_count_and_media_type(deck: DeckSpec, options: RenderOptions) -> None:
    rendered = render_presentation(deck, options)
    assert rendered.slide_count == len(deck.slides)
    assert rendered.media_type == MEDIA_TYPES["pptx"]
    assert rendered.filename.endswith(".pptx")
    assert rendered.size_bytes > 0


def test_output_is_a_real_ooxml_package(deck: DeckSpec, options: RenderOptions) -> None:
    # A .pptx is a zip; the magic bytes catch a renderer that returns something else
    # entirely, which python-pptx would refuse far less clearly.
    assert render_presentation(deck, options).data[:2] == b"PK"


# --------------------------------------------------------------------------- #
# Content
# --------------------------------------------------------------------------- #


def test_nested_bullets_keep_their_indent_level(options: RenderOptions) -> None:
    spec = DeckSpec(
        title="T",
        slides=[
            Slide(
                layout="title_content",
                title="Ebenen",
                bullets=[
                    Bullet(
                        text="eins",
                        children=[Bullet(text="zwei", children=[Bullet(text="drei")])],
                    )
                ],
            )
        ],
    )
    slide = opened(render_presentation(spec, options)).slides[0]
    body = next(
        shape
        for shape in slide.shapes
        if shape.has_text_frame and "eins" in shape.text_frame.text
    )
    levels = {p.text: p.level for p in body.text_frame.paragraphs if p.text}
    assert levels == {"eins": 0, "zwei": 1, "drei": 2}


def test_table_slide_carries_headers_and_rows(options: RenderOptions) -> None:
    spec = DeckSpec(
        title="T",
        slides=[
            Slide(
                layout="table",
                title="Regionen",
                table=TableData(
                    headers=["Region", "Umsatz"],
                    rows=[["DACH", "4.2M"], ["UK", "1.1M"]],
                ),
            )
        ],
    )
    slide = opened(render_presentation(spec, options)).slides[0]
    table = next(shape.table for shape in slide.shapes if shape.has_table)
    assert [cell.text for cell in table.rows[0].cells] == ["Region", "Umsatz"]
    assert len(table.rows) == 3  # header + two data rows


def test_chart_slide_produces_a_real_chart(options: RenderOptions) -> None:
    spec = DeckSpec(
        title="T",
        slides=[
            Slide(
                layout="chart",
                title="Verlauf",
                chart=ChartData(
                    categories=["Q1", "Q2"],
                    series=[ChartSeries(name="2026", values=[1.0, 2.0])],
                ),
            )
        ],
    )
    slide = opened(render_presentation(spec, options)).slides[0]
    charts = [shape.chart for shape in slide.shapes if shape.has_chart]
    assert len(charts) == 1
    assert [s.name for s in charts[0].plots[0].series] == ["2026"]


def test_inline_image_is_embedded(tiny_png: str, options: RenderOptions) -> None:
    spec = DeckSpec(
        title="T",
        slides=[
            Slide(
                layout="image",
                title="Bild",
                image=ImageRef(data_base64=tiny_png, width_cm=6),
            )
        ],
    )
    slide = opened(render_presentation(spec, options)).slides[0]
    pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
    assert len(pictures) == 1
    assert pictures[0].width == pytest.approx(Emu(6 * 360000), rel=0.01)


def test_alt_text_reaches_the_shape(tiny_png: str, options: RenderOptions) -> None:
    spec = DeckSpec(
        title="T",
        slides=[
            Slide(
                layout="image",
                image=ImageRef(data_base64=tiny_png, alt_text="Ein Diagramm"),
            )
        ],
    )
    slide = opened(render_presentation(spec, options)).slides[0]
    picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
    assert "Ein Diagramm" in picture._element.nvPicPr.cNvPr.get("descr", "")


# --------------------------------------------------------------------------- #
# Speaker notes
# --------------------------------------------------------------------------- #


def test_notes_are_written_when_requested() -> None:
    spec = DeckSpec(
        title="T",
        slides=[Slide(layout="title_content", title="A", notes="Langsam sprechen.")],
    )
    rendered = render_presentation(spec, RenderOptions(include_notes=True))
    slide = opened(rendered).slides[0]
    assert slide.has_notes_slide
    assert "Langsam sprechen." in slide.notes_slide.notes_text_frame.text


def test_include_notes_does_not_suppress_notes_the_caller_supplied() -> None:
    """``include_notes`` gates *generation*, not explicit content.

    It is a hint for brief mode — "invent speaker notes" — so a spec that already carries
    notes keeps them either way. Dropping content a caller explicitly wrote because a
    formatting flag was off would be the wrong reading of the option.
    """
    spec = DeckSpec(
        title="T",
        slides=[Slide(layout="title_content", title="A", notes="Trotzdem zeigen.")],
    )
    rendered = render_presentation(spec, RenderOptions(include_notes=False))
    slide = opened(rendered).slides[0]
    assert slide.has_notes_slide
    assert "Trotzdem zeigen." in slide.notes_slide.notes_text_frame.text


def test_a_slide_without_notes_gets_no_notes_slide() -> None:
    spec = DeckSpec(title="T", slides=[Slide(layout="title_content", title="A")])
    slide = opened(render_presentation(spec, RenderOptions(include_notes=True))).slides[0]
    assert not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip()


# --------------------------------------------------------------------------- #
# Errors
# --------------------------------------------------------------------------- #


def test_render_error_names_the_slide_that_failed(options: RenderOptions) -> None:
    """The position is the whole point of the message.

    A bare re-raise once lost it, leaving the model with "image could not be decoded" and
    no way to know which of twenty slides to fix.
    """
    spec = DeckSpec(
        title="T",
        slides=[
            Slide(layout="title", title="ok"),
            Slide(layout="image", image=ImageRef(data_base64="not-base64-at-all")),
        ],
    )
    with pytest.raises(RenderError) as caught:
        render_presentation(spec, options)
    assert "2" in str(caught.value)


def test_unresolvable_file_id_fails_without_a_resolver(options: RenderOptions) -> None:
    spec = DeckSpec(
        title="T",
        slides=[Slide(layout="image", image=ImageRef(file_id="f-404"))],
    )
    with pytest.raises(RenderError):
        render_presentation(spec, options, image_resolver=None)


# --------------------------------------------------------------------------- #
# Options
# --------------------------------------------------------------------------- #


def test_unknown_font_is_warned_about_but_still_renders(deck: DeckSpec) -> None:
    rendered = render_presentation(deck, RenderOptions(font_family="Nicht Installiert"))
    assert rendered.size_bytes > 0
    assert any("Nicht Installiert" in warning for warning in rendered.warnings)


def test_safe_font_produces_no_font_warning(deck: DeckSpec) -> None:
    rendered = render_presentation(deck, RenderOptions(font_family="Arial"))
    assert not any("Arial" in warning for warning in rendered.warnings)


def test_filename_option_is_honoured(deck: DeckSpec) -> None:
    rendered = render_presentation(deck, RenderOptions(filename="Quartal Q3"))
    assert rendered.filename.startswith("Quartal")
    assert rendered.filename.endswith(".pptx")


def test_filename_cannot_escape_its_directory(deck: DeckSpec) -> None:
    rendered = render_presentation(deck, RenderOptions(filename="../../etc/passwd"))
    assert "/" not in rendered.filename
    assert ".." not in rendered.filename
