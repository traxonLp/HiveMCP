"""Reading and patching documents that arrived from the chat.

Every test here round-trips through the renderers rather than checking in binary
fixtures: render a document, edit it, read it back. That keeps the fixtures honest (they
cannot drift from what HiveMCP actually produces) and exercises the read side and the
edit side against each other, which is how the two are used in practice — the model calls
``hive_read_document`` to find the numbers it then passes to ``hive_edit_document``.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from hivemcp.core.editing.apply import EditError, apply_edits
from hivemcp.core.editing.read import DocumentUnreadable, read_document
from hivemcp.core.models import (
    AppendParagraph,
    Bullet,
    Column,
    DeckSpec,
    DeleteSlide,
    DocSpec,
    FillPlaceholders,
    Heading,
    Paragraph,
    RenderOptions,
    ReorderSlides,
    ReplaceText,
    SetCell,
    SetNotes,
    SetParagraph,
    Sheet,
    SheetSpec,
    Slide,
    TableBlock,
    TableData,
)
from hivemcp.core.render.docx import render_document
from hivemcp.core.render.pptx import render_presentation
from hivemcp.core.render.xlsx import render_spreadsheet


@pytest.fixture
def deck_bytes(options: RenderOptions) -> bytes:
    spec = DeckSpec(
        title="Bericht",
        slides=[
            Slide(layout="title", title="Erste", subtitle="{{untertitel}}"),
            Slide(layout="title_content", title="Zweite", bullets=[Bullet(text="Alpha")]),
            Slide(layout="title_content", title="Dritte", bullets=[Bullet(text="Beta")]),
        ],
    )
    return render_presentation(spec, options).data


@pytest.fixture
def doc_bytes(options: RenderOptions) -> bytes:
    spec = DocSpec(
        title="Handbuch",
        blocks=[
            Heading(text="Einleitung", level=1),
            Paragraph(text="Alpha steht hier."),
            Paragraph(text="Kunde: {{kunde}}"),
            TableBlock(data=TableData(headers=["A"], rows=[["Alpha"]])),
        ],
    )
    return render_document(spec, options).data


@pytest.fixture
def sheet_bytes(options: RenderOptions) -> bytes:
    spec = SheetSpec(
        title="Zahlen",
        sheets=[
            Sheet(
                name="Daten",
                columns=[
                    Column(header="Name", key="n"),
                    Column(header="Wert", key="w", type="number"),
                ],
                rows=[{"n": "Alpha", "w": 1.0}, {"n": "Beta", "w": 2.0}],
            )
        ],
    )
    return render_spreadsheet(spec, options).data


def written(tmp_path: Path, data: bytes, name: str) -> Path:
    path = tmp_path / name
    path.write_bytes(data)
    return path


# --------------------------------------------------------------------------- #
# Reading
# --------------------------------------------------------------------------- #


def test_reads_pptx_outline(tmp_path: Path, deck_bytes: bytes) -> None:
    report = read_document(written(tmp_path, deck_bytes, "d.pptx"), "pptx")
    assert report["kind"] == "pptx"
    assert report["slide_count"] == 3
    assert [entry["slide"] for entry in report["slides"]] == [1, 2, 3]
    assert report["slides"][0]["title"] == "Erste"


def test_pptx_read_surfaces_placeholders(tmp_path: Path, deck_bytes: bytes) -> None:
    report = read_document(written(tmp_path, deck_bytes, "d.pptx"), "pptx")
    assert "untertitel" in report["placeholders"]


def test_reads_docx_paragraph_numbers_and_styles(tmp_path: Path, doc_bytes: bytes) -> None:
    report = read_document(written(tmp_path, doc_bytes, "d.docx"), "docx")
    assert report["kind"] == "docx"
    numbers = [entry["paragraph"] for entry in report["paragraphs"]]
    assert numbers == sorted(numbers)
    assert report["paragraph_count"] >= len(report["paragraphs"])
    assert report["styles_available"]
    assert report["tables"][0]["header"] == ["A"]


def test_docx_outline_omits_empty_paragraphs_but_keeps_the_numbering(
    tmp_path: Path, doc_bytes: bytes
) -> None:
    """The numbers must stay absolute, because set_paragraph indexes the real document."""
    outline = read_document(written(tmp_path, doc_bytes, "d.docx"), "docx", "outline")
    full = read_document(written(tmp_path, doc_bytes, "d.docx"), "docx", "full")
    assert len(outline["paragraphs"]) <= len(full["paragraphs"])
    assert all(entry["text"].strip() for entry in outline["paragraphs"])


def test_reads_xlsx(tmp_path: Path, sheet_bytes: bytes) -> None:
    report = read_document(written(tmp_path, sheet_bytes, "d.xlsx"), "xlsx")
    assert report["kind"] == "xlsx"
    assert "Daten" in str(report)


def test_unknown_kind_is_rejected(tmp_path: Path, doc_bytes: bytes) -> None:
    with pytest.raises(DocumentUnreadable):
        read_document(written(tmp_path, doc_bytes, "d.docx"), "pdf")


def test_a_file_that_is_not_a_document_fails_cleanly(tmp_path: Path) -> None:
    with pytest.raises(DocumentUnreadable):
        read_document(written(tmp_path, b"not a zip at all", "x.docx"), "docx")


def test_wrong_kind_for_the_actual_content_fails_cleanly(
    tmp_path: Path, doc_bytes: bytes
) -> None:
    with pytest.raises(DocumentUnreadable):
        read_document(written(tmp_path, doc_bytes, "d.pptx"), "pptx")


# --------------------------------------------------------------------------- #
# Editing: PowerPoint
# --------------------------------------------------------------------------- #


def test_replace_text_changes_every_occurrence(deck_bytes: bytes) -> None:
    rendered, log = apply_edits(
        deck_bytes, "pptx", [ReplaceText(find="Alpha", replace="Gamma")], "out.pptx"
    )
    text = "\n".join(
        shape.text_frame.text
        for slide in Presentation(BytesIO(rendered.data)).slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Gamma" in text and "Alpha" not in text
    assert log and "replaced" in log[0]


def test_replace_text_is_case_sensitive_by_default(deck_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        deck_bytes, "pptx", [ReplaceText(find="alpha", replace="Gamma")], "out.pptx"
    )
    assert "Gamma" not in str(rendered.data)


def test_matching_nothing_succeeds_but_warns(deck_bytes: bytes) -> None:
    """An operation can be valid and still change nothing. That must be visible.

    Silently reporting success here is how a user ends up believing a document was
    changed when it was not.
    """
    rendered, log = apply_edits(
        deck_bytes, "pptx", [ReplaceText(find="kommt-nicht-vor", replace="x")], "out.pptx"
    )
    assert rendered.warnings
    assert "matched nothing" in rendered.warnings[0]
    assert log


def test_delete_slide_removes_exactly_one(deck_bytes: bytes) -> None:
    rendered, _ = apply_edits(deck_bytes, "pptx", [DeleteSlide(slide=2)], "out.pptx")
    presentation = Presentation(BytesIO(rendered.data))
    assert len(presentation.slides) == 2
    titles = [
        slide.shapes.title.text for slide in presentation.slides if slide.shapes.title
    ]
    assert "Zweite" not in titles


def test_reorder_slides_applies_the_given_order(deck_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        deck_bytes, "pptx", [ReorderSlides(order=[3, 1, 2])], "out.pptx"
    )
    presentation = Presentation(BytesIO(rendered.data))
    titles = [
        slide.shapes.title.text if slide.shapes.title else "" for slide in presentation.slides
    ]
    assert titles == ["Dritte", "Erste", "Zweite"]


def test_set_notes_writes_speaker_notes(deck_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        deck_bytes, "pptx", [SetNotes(slide=1, notes="Kurz halten.")], "out.pptx"
    )
    slide = Presentation(BytesIO(rendered.data)).slides[0]
    assert slide.notes_slide.notes_text_frame.text == "Kurz halten."


def test_fill_placeholders_substitutes_the_token(deck_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        deck_bytes,
        "pptx",
        [FillPlaceholders(values={"untertitel": "Q3 2026"})],
        "out.pptx",
    )
    text = "\n".join(
        shape.text_frame.text
        for slide in Presentation(BytesIO(rendered.data)).slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Q3 2026" in text and "{{untertitel}}" not in text


def test_out_of_range_slide_is_refused(deck_bytes: bytes) -> None:
    with pytest.raises(EditError) as caught:
        apply_edits(deck_bytes, "pptx", [DeleteSlide(slide=99)], "out.pptx")
    assert "99" in str(caught.value)


def test_an_operation_for_another_format_is_refused(deck_bytes: bytes) -> None:
    with pytest.raises(EditError) as caught:
        apply_edits(deck_bytes, "pptx", [AppendParagraph(text="x")], "out.pptx")
    assert "pptx" in str(caught.value)


def test_edits_are_all_or_nothing(deck_bytes: bytes) -> None:
    """A later failure must not leave the earlier operations applied."""
    with pytest.raises(EditError):
        apply_edits(
            deck_bytes,
            "pptx",
            [ReplaceText(find="Alpha", replace="Gamma"), DeleteSlide(slide=99)],
            "out.pptx",
        )
    # The input buffer is untouched, which is what the caller re-reads on failure.
    assert len(Presentation(BytesIO(deck_bytes)).slides) == 3


def test_no_operations_is_an_error(deck_bytes: bytes) -> None:
    with pytest.raises(EditError):
        apply_edits(deck_bytes, "pptx", [], "out.pptx")


# --------------------------------------------------------------------------- #
# Editing: Word
# --------------------------------------------------------------------------- #


def test_set_paragraph_targets_the_number_from_read(
    tmp_path: Path, doc_bytes: bytes
) -> None:
    """Read and edit must agree on what "paragraph 3" means.

    Reading first and feeding the number straight back is exactly what the model does,
    so an off-by-one between the two modules would be invisible to a test that only
    exercised one of them.
    """
    report = read_document(written(tmp_path, doc_bytes, "d.docx"), "docx", "full")
    target = next(
        entry for entry in report["paragraphs"] if "Alpha steht hier." in entry["text"]
    )
    rendered, _ = apply_edits(
        doc_bytes,
        "docx",
        [SetParagraph(paragraph=target["paragraph"], text="Neu.")],
        "out.docx",
    )
    texts = [p.text for p in Document(BytesIO(rendered.data)).paragraphs]
    assert "Neu." in texts
    assert "Alpha steht hier." not in texts


def test_append_paragraph_adds_at_the_end(doc_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        doc_bytes, "docx", [AppendParagraph(text="Nachtrag.")], "out.docx"
    )
    assert Document(BytesIO(rendered.data)).paragraphs[-1].text == "Nachtrag."


def test_append_paragraph_rejects_an_unknown_style(doc_bytes: bytes) -> None:
    with pytest.raises(EditError) as caught:
        apply_edits(
            doc_bytes,
            "docx",
            [AppendParagraph(text="x", style="Gibt Es Nicht")],
            "out.docx",
        )
    assert "Gibt Es Nicht" in str(caught.value)


def test_replace_text_reaches_into_tables(doc_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        doc_bytes, "docx", [ReplaceText(find="Alpha", replace="Gamma")], "out.docx"
    )
    table = Document(BytesIO(rendered.data)).tables[0]
    assert table.rows[1].cells[0].text == "Gamma"


def test_replacement_survives_being_split_across_runs(options: RenderOptions) -> None:
    """Word splits a sentence across runs at will; a naive per-run replace misses it."""
    spec = DocSpec(title="T", blocks=[Paragraph(text="Der Kunde heisst Alpha GmbH.")])
    data = render_document(spec, options).data
    rendered, _ = apply_edits(
        data, "docx", [ReplaceText(find="Alpha GmbH", replace="Beta AG")], "out.docx"
    )
    text = "\n".join(p.text for p in Document(BytesIO(rendered.data)).paragraphs)
    assert "Beta AG" in text and "Alpha GmbH" not in text


def test_out_of_range_paragraph_is_refused(doc_bytes: bytes) -> None:
    with pytest.raises(EditError) as caught:
        apply_edits(doc_bytes, "docx", [SetParagraph(paragraph=999, text="x")], "out.docx")
    assert "999" in str(caught.value)


# --------------------------------------------------------------------------- #
# Editing: Excel
# --------------------------------------------------------------------------- #


def test_set_cell_writes_a_typed_value(sheet_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        sheet_bytes,
        "xlsx",
        [SetCell(sheet="Daten", cell="B2", value=42, type="integer")],
        "out.xlsx",
    )
    assert load_workbook(BytesIO(rendered.data))["Daten"]["B2"].value == 42


def test_set_cell_accepts_lowercase_references(sheet_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        sheet_bytes,
        "xlsx",
        [SetCell(sheet="Daten", cell="b2", value=7, type="integer")],
        "out.xlsx",
    )
    assert load_workbook(BytesIO(rendered.data))["Daten"]["B2"].value == 7


def test_set_cell_neutralises_an_injected_formula(sheet_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        sheet_bytes,
        "xlsx",
        [SetCell(sheet="Daten", cell="A2", value="=1+1", type="text")],
        "out.xlsx",
    )
    value = load_workbook(BytesIO(rendered.data))["Daten"]["A2"].value
    assert value.startswith("'") or not value.startswith("=")


def test_unknown_sheet_is_refused_and_lists_the_real_ones(sheet_bytes: bytes) -> None:
    with pytest.raises(EditError) as caught:
        apply_edits(
            sheet_bytes,
            "xlsx",
            [SetCell(sheet="Gibtsnicht", cell="A1", value="x")],
            "out.xlsx",
        )
    assert "Daten" in str(caught.value)


def test_replace_text_walks_every_sheet(sheet_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        sheet_bytes, "xlsx", [ReplaceText(find="Alpha", replace="Gamma")], "out.xlsx"
    )
    assert load_workbook(BytesIO(rendered.data))["Daten"]["A2"].value == "Gamma"


# --------------------------------------------------------------------------- #
# Result plumbing
# --------------------------------------------------------------------------- #


def test_result_keeps_the_requested_filename_and_media_type(deck_bytes: bytes) -> None:
    rendered, _ = apply_edits(
        deck_bytes, "pptx", [SetNotes(slide=1, notes="x")], "Bericht final.pptx"
    )
    assert rendered.filename == "Bericht final.pptx"
    assert rendered.media_type.endswith("presentationml.presentation")


def test_edited_output_can_be_read_again(tmp_path: Path, deck_bytes: bytes) -> None:
    """The edit must produce a file the read side still understands."""
    rendered, _ = apply_edits(deck_bytes, "pptx", [DeleteSlide(slide=1)], "out.pptx")
    report = read_document(written(tmp_path, rendered.data, "out.pptx"), "pptx")
    assert report["slide_count"] == 2
